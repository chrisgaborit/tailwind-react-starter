import os
import json
from tqdm import tqdm
from PyPDF2 import PdfReader
from pptx import Presentation
import google.generativeai as genai

# Set your Google API Key (Gemini API Key)
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY") or "YOUR_GEMINI_API_KEY"
genai.configure(api_key=GOOGLE_API_KEY)

CONVERSION_PROMPT = """
You are an expert eLearning instructional designer and developer.
Given any human-created storyboard (text, PPT, PDF, etc.), convert it into a production-ready JSON blueprint for an eLearning module.

Your JSON output must include the following for every scene:
- sceneNumber, sceneTitle
- voiceover
- onScreenText (array: separate heading, body, instructions, labels, etc.)
- layout: detailed instructions for visual and text layout (e.g., image left, text right, CTA bottom)
- mediaAssets: array of objects (type: image/video/icon, filename, altText, fileUrl if known)
- interactivity: type (e.g., click-to-reveal, MCQ, drag-and-drop), full interaction logic (button labels, options, outcomes, nextScene or branching, state changes)
- knowledgeCheck: question, options, correctAnswer(s), feedback, remediation (nextScene if wrong/correct), scoring if relevant
- branching/navigation: nextScene (by number or key), branchingOptions (for scenarios or choices)
- accessibility: altText for every visual, transcript for every audio/voiceover, ARIA label or keyboard notes if present
- animation/timing: timing, animation type, trigger if specified
- coachGuidance: guidance for facilitators or AI coaches
- authorNotes: anything for developers or content creators

Global/module-level fields must include:
- moduleName, moduleType, duration, complexityLevel, tags, learningOutcomes, audience, brandGuidelines (colours, fonts), sponsorLogo
- glossary, references, or extra fields if present

If information is not in the source, output an empty string, array, or "TBD" as placeholder—do not invent content.

Format the final output as prettified JSON only, with no extra explanation or comments.
"""

def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    return "\n".join([page.extract_text() or '' for page in reader.pages])

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def convert_to_json_with_gemini(storyboard_text):
    model = genai.GenerativeModel('gemini-1.5-flash')
    response = model.generate_content([
        CONVERSION_PROMPT,
        storyboard_text[:25000]  # Gemini Flash handles more context
    ])
    # The response is in response.text
    json_response = response.text
    try:
        return json.loads(json_response)
    except Exception:
        json_response = json_response.strip('` \n')
        return json.loads(json_response)

def process_folder(input_folder, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    for filename in tqdm(os.listdir(input_folder)):
        if filename.endswith(".pdf"):
            text = extract_text_from_pdf(os.path.join(input_folder, filename))
        elif filename.endswith(".pptx"):
            text = extract_text_from_pptx(os.path.join(input_folder, filename))
        elif filename.endswith(".txt"):
            with open(os.path.join(input_folder, filename), "r", encoding="utf-8") as f:
                text = f.read()
        else:
            continue

        try:
            json_data = convert_to_json_with_gemini(text)
            base_name = os.path.splitext(filename)[0]
            with open(os.path.join(output_folder, f"{base_name}.json"), "w", encoding="utf-8") as out_f:
                json.dump(json_data, out_f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"Failed to process {filename}: {e}")

if __name__ == "__main__":
    input_folder = "storyboards_input"
    output_folder = "storyboards_json"
    process_folder(input_folder, output_folder)
